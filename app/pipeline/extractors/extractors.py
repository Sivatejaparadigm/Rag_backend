from __future__ import annotations

import csv
import re
import uuid
from pathlib import Path

from app.pipeline.extractors.base import BaseExtractor, ExtractionError
from app.schemas.ingestion import DocumentType, ExtractedContentCreate


# ── PDF ───────────────────────────────────────────────────────

class PDFExtractor(BaseExtractor):
    supported_type = DocumentType.PDF

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            from pypdf import PdfReader
        except ImportError:
            raise ExtractionError("pypdf not installed. Run: pip install pypdf")

        try:
            self._ensure_file_exists(file_path)
            reader = PdfReader(str(file_path))
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to read PDF", original=e) from e

        pages = []
        all_text = []

        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            all_text.append(text)
            pages.append({
                "page_number": i,
                "text": text,
                "word_count": self._word_count(text),
            })

        raw_text = self._truncate("\n\n".join(filter(None, all_text)))

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=raw_text,
            pages=pages,
            tables=[],
            warnings=[],
        )


# ── DOCX ──────────────────────────────────────────────────────

class DocxExtractor(BaseExtractor):
    supported_type = DocumentType.DOCX

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            from docx import Document
        except ImportError:
            raise ExtractionError("python-docx not installed. Run: pip install python-docx")

        try:
            self._ensure_file_exists(file_path)
            doc = Document(str(file_path))
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to open DOCX file", original=e) from e

        paragraphs = []
        all_text = []
        tables = []

        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if not text:
                continue
            all_text.append(text)
            paragraphs.append({
                "index": i,
                "style": para.style.name if para.style else "Normal",
                "is_heading": para.style.name.startswith("Heading") if para.style else False,
                "text": text,
                "word_count": self._word_count(text),
            })

        for i, table in enumerate(doc.tables):
            rows = [
                [cell.text.strip() for cell in row.cells]
                for row in table.rows
            ]
            if rows:
                tables.append({
                    "index": i,
                    "headers": rows[0],
                    "rows": rows[1:],
                })

        raw_text = self._truncate("\n\n".join(all_text))

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=raw_text,
            pages=paragraphs,
            tables=tables,
            warnings=[],
        )


# ── TXT ───────────────────────────────────────────────────────

class TextExtractor(BaseExtractor):
    supported_type = DocumentType.TXT

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            self._ensure_file_exists(file_path)
            text = self._truncate(file_path.read_bytes().decode("utf-8", errors="ignore"))
        except Exception as e:
            raise ExtractionError("Failed to read text file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=text,
            pages=[{"page_number": 1, "text": text, "word_count": self._word_count(text)}],
            tables=[],
            warnings=[],
        )


# ── CSV ───────────────────────────────────────────────────────

class CSVExtractor(BaseExtractor):
    supported_type = DocumentType.CSV

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            self._ensure_file_exists(file_path)
            text = file_path.read_bytes().decode("utf-8", errors="ignore")
            rows = list(csv.reader(text.splitlines()))
        except Exception as e:
            raise ExtractionError("Failed to read CSV file", original=e) from e

        headers = rows[0] if rows else []
        data_rows = rows[1:] if len(rows) > 1 else []
        raw_text = self._truncate("\n".join("\t".join(r) for r in rows))

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=raw_text,
            pages=[{"page_number": 1, "text": raw_text, "word_count": self._word_count(raw_text)}],
            tables=[{"headers": headers, "rows": data_rows}],
            warnings=[],
        )


# ── HTML ──────────────────────────────────────────────────────

class HTMLExtractor(BaseExtractor):
    supported_type = DocumentType.HTML

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ExtractionError("beautifulsoup4 not installed. Run: pip install beautifulsoup4 lxml")

        try:
            self._ensure_file_exists(file_path)
            soup = BeautifulSoup(file_path.read_text(encoding="utf-8", errors="ignore"), "lxml")
            for tag in soup(["script", "style", "nav", "footer"]):
                tag.decompose()
            raw_text = self._truncate(soup.get_text(separator="\n").strip())
            tables = []
            for i, table in enumerate(soup.find_all("table")):
                rows = [
                    [td.get_text(strip=True) for td in row.find_all(["th", "td"])]
                    for row in table.find_all("tr")
                ]
                if rows:
                    tables.append({"index": i, "headers": rows[0], "rows": rows[1:]})
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to read HTML file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=raw_text,
            pages=[{"page_number": 1, "text": raw_text, "word_count": self._word_count(raw_text)}],
            tables=tables,
            warnings=[],
        )


# ── Markdown ──────────────────────────────────────────────────

class MarkdownExtractor(BaseExtractor):
    supported_type = DocumentType.MD

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            self._ensure_file_exists(file_path)
            text = file_path.read_text(encoding="utf-8", errors="ignore")
            plain = re.sub(r"#{1,6}\s", "", text)
            plain = re.sub(r"\*{1,2}(.+?)\*{1,2}", r"\1", plain)
            plain = re.sub(r"`{1,3}.+?`{1,3}", "", plain, flags=re.DOTALL)
            plain = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", plain)
            plain = self._truncate(plain.strip())
        except Exception as e:
            raise ExtractionError("Failed to read Markdown file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=plain,
            pages=[{"page_number": 1, "text": plain, "word_count": self._word_count(plain)}],
            tables=[],
            warnings=[],
        )


# ── PPTX ──────────────────────────────────────────────────────

class PPTXExtractor(BaseExtractor):
    supported_type = DocumentType.PPTX

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            from pptx import Presentation
        except ImportError:
            raise ExtractionError("python-pptx not installed. Run: pip install python-pptx")

        try:
            self._ensure_file_exists(file_path)
            prs = Presentation(str(file_path))
            slides = []
            all_text = []
            for i, slide in enumerate(prs.slides, start=1):
                text = "\n".join(
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                )
                all_text.append(text)
                slides.append({
                    "page_number": i,
                    "text": text,
                    "word_count": self._word_count(text),
                })
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to read PPTX file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=self._truncate("\n\n".join(filter(None, all_text))),
            pages=slides,
            tables=[],
            warnings=[],
        )


# ── XLSX ──────────────────────────────────────────────────────

class XLSXExtractor(BaseExtractor):
    supported_type = DocumentType.XLSX

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            import openpyxl
        except ImportError:
            raise ExtractionError("openpyxl not installed. Run: pip install openpyxl")

        try:
            self._ensure_file_exists(file_path)
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            sheets = []
            tables = []
            all_text = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue
                headers = [str(c) if c is not None else "" for c in rows[0]]
                data_rows = [
                    [str(c) if c is not None else "" for c in row]
                    for row in rows[1:]
                ]
                text = "\n".join("\t".join(r) for r in ([headers] + data_rows))
                all_text.append(f"[{sheet_name}]\n{text}")
                sheets.append({
                    "sheet": sheet_name,
                    "text": text,
                    "word_count": self._word_count(text),
                })
                tables.append({
                    "sheet": sheet_name,
                    "headers": headers,
                    "rows": data_rows,
                })
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to read XLSX file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=self._truncate("\n\n".join(all_text)),
            pages=sheets,
            tables=tables,
            warnings=[],
        )


# ── RTF ───────────────────────────────────────────────────────

class RTFExtractor(BaseExtractor):
    supported_type = DocumentType.RTF

    async def extract(self, file_path: Path, tenant_id: uuid.UUID) -> ExtractedContentCreate:
        try:
            from striprtf.striprtf import rtf_to_text
        except ImportError:
            raise ExtractionError("striprtf not installed. Run: pip install striprtf")

        try:
            self._ensure_file_exists(file_path)
            raw_text = self._truncate(
                rtf_to_text(
                    file_path.read_text(encoding="utf-8", errors="ignore")
                ).strip()
            )
        except ExtractionError:
            raise
        except Exception as e:
            raise ExtractionError("Failed to read RTF file", original=e) from e

        return ExtractedContentCreate(
            tenant_id=tenant_id,
            raw_text=raw_text,
            pages=[{"page_number": 1, "text": raw_text, "word_count": self._word_count(raw_text)}],
            tables=[],
            warnings=[],
        )